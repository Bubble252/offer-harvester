from __future__ import annotations

import difflib
import hashlib
import re
from pathlib import Path
from typing import Iterable, List, Literal, Optional

from models import AdvisorProfile, RAGSearchHit, StudentProfile, now_iso
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import BaseModel, Field

SlideElementKind = Literal[
    "title", "body", "image", "logo", "date", "footer", "shape", "unsupported"
]


class SlideElementSchema(BaseModel):
    element_id: str
    kind: SlideElementKind
    description: str = ""
    char_capacity: int = 0
    image_slot: bool = False
    constraints: List[str] = Field(default_factory=list)


class SlideSchema(BaseModel):
    slide_index: int
    page_type: str = "content"
    title: str = ""
    elements: List[SlideElementSchema] = Field(default_factory=list)
    warnings: List[str] = Field(default_factory=list)


class PresentationSchema(BaseModel):
    schema_id: str
    source_path: str = ""
    source_hash: str = ""
    schema_version: str = "17c-v1"
    slides: List[SlideSchema] = Field(default_factory=list)
    skipped_elements: List[str] = Field(default_factory=list)
    warnings: List[str] = Field(default_factory=list)
    created_at: str = Field(default_factory=now_iso)


class OrganizedSlideContent(BaseModel):
    slide_index: int
    paragraph: str = ""
    bullets: List[str] = Field(default_factory=list)
    speaker_notes: List[str] = Field(default_factory=list)
    evidence_refs: List[str] = Field(default_factory=list)
    unconfirmed_fields: List[str] = Field(default_factory=list)
    rejected_fields: List[str] = Field(default_factory=list)
    character_budget: int = 0


class ContentOrganizationResult(BaseModel):
    organizer_version: str = "17d-v1"
    slides: List[OrganizedSlideContent] = Field(default_factory=list)
    evidence_refs: List[str] = Field(default_factory=list)
    warnings: List[str] = Field(default_factory=list)
    created_at: str = Field(default_factory=now_iso)


class SingleSlideCopilotRequest(BaseModel):
    slide_index: int
    instruction: str
    original: OrganizedSlideContent


class SingleSlideCopilotResult(BaseModel):
    copilot_version: str = "17e-v1"
    slide_index: int
    original: OrganizedSlideContent
    updated: OrganizedSlideContent
    diff_text: str = ""
    passed: bool = True
    issues: List[str] = Field(default_factory=list)
    action_items: List[str] = Field(default_factory=list)
    created_at: str = Field(default_factory=now_iso)


class SingleSlideVersionRecord(BaseModel):
    version_id: str
    slide_index: int
    instruction: str
    original: OrganizedSlideContent
    updated: OrganizedSlideContent
    diff_text: str = ""
    passed: bool = True
    issues: List[str] = Field(default_factory=list)
    action_items: List[str] = Field(default_factory=list)
    created_at: str = Field(default_factory=now_iso)


def build_presentation_schema(path: Path, *, source_reference_id: str = "") -> PresentationSchema:
    """Extract a conservative, auditable schema without attempting visual reconstruction."""

    source_hash = f"sha256:{hashlib.sha256(path.read_bytes()).hexdigest()}"
    presentation = Presentation(str(path))
    slides: List[SlideSchema] = []
    skipped: List[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = [shape.text.strip() for shape in slide.shapes if getattr(shape, "text", "").strip()]
        title = texts[0] if texts else ""
        elements: List[SlideElementSchema] = []
        warnings: List[str] = []
        text_seen = 0
        for shape_index, shape in enumerate(slide.shapes, start=1):
            element_id = f"slide-{slide_index}-element-{shape_index}"
            kind = _element_kind(shape, text_seen)
            if kind in {"title", "body", "date", "footer"}:
                text_seen += 1
            if kind == "unsupported":
                skipped.append(element_id)
                warnings.append(f"{element_id}: unsupported shape type")
            elements.append(
                SlideElementSchema(
                    element_id=element_id,
                    kind=kind,
                    description=_shape_description(shape),
                    char_capacity=_char_capacity(shape),
                    image_slot=kind in {"image", "logo"},
                    constraints=_constraints(kind),
                )
            )
        slides.append(
            SlideSchema(
                slide_index=slide_index,
                page_type=_page_type(slide_index, texts),
                title=title[:120],
                elements=elements,
                warnings=warnings,
            )
        )
    return PresentationSchema(
        schema_id=f"schema_{source_hash.split(':', 1)[1][:12]}",
        source_path=source_reference_id or path.name,
        source_hash=source_hash,
        slides=slides,
        skipped_elements=skipped,
        warnings=["复杂形状只保留 unsupported 占位，不静默丢弃。"] if skipped else [],
    )


def organize_profile_content(
    profile: StudentProfile,
    advisor: Optional[AdvisorProfile] = None,
    evidence_hits: Optional[Iterable[RAGSearchHit]] = None,
    *,
    slide_index: int = 1,
    character_budget: int = 520,
) -> ContentOrganizationResult:
    """Turn existing facts into presentation-friendly forms; never invent new facts."""

    evidence_hits = list(evidence_hits or [])
    evidence_refs = [hit.evidence_ref for hit in evidence_hits if hit.evidence_ref]
    rejected = {field for field, status in profile.confirmation_map.items() if status == "rejected"}
    unconfirmed = {
        field
        for field, status in profile.confirmation_map.items()
        if status in {"unconfirmed", "needs_review"}
    }
    projects = [] if "projects" in rejected else list(profile.projects)
    interests = [] if "research_interests" in rejected else list(profile.research_interests)
    advisor_directions = (
        list(advisor.research_directions[:3]) if advisor and advisor.research_directions else []
    )
    bullets = [f"项目：{item}" for item in projects[:3]]
    bullets.extend(f"兴趣：{item}" for item in interests[:2])
    if advisor_directions:
        bullets.append(f"导师方向：{'、'.join(advisor_directions)}")
    bullets = [item[:140] for item in bullets]
    paragraph_parts = []
    if profile.education and "education" not in rejected:
        paragraph_parts.append(profile.education)
    if projects:
        paragraph_parts.append(f"主要项目包括{'、'.join(projects[:2])}")
    if interests:
        paragraph_parts.append(f"研究兴趣集中在{'、'.join(interests[:2])}")
    paragraph = "。".join(paragraph_parts)
    if paragraph and not paragraph.endswith(("。", "！", "？")):
        paragraph += "。"
    notes = ["页面内容由现有 profile、导师方向和 evidence 组织，不新增事实。"]
    if unconfirmed:
        notes.append(f"以下字段仍需确认：{', '.join(sorted(unconfirmed))}")
    slide = OrganizedSlideContent(
        slide_index=slide_index,
        paragraph=paragraph[:character_budget],
        bullets=bullets,
        speaker_notes=notes,
        evidence_refs=evidence_refs,
        unconfirmed_fields=sorted(unconfirmed),
        rejected_fields=sorted(rejected),
        character_budget=character_budget,
    )
    warnings = ["存在未确认字段，提交前需要人工确认。"] if unconfirmed else []
    return ContentOrganizationResult(
        slides=[slide],
        evidence_refs=evidence_refs,
        warnings=warnings,
    )


def run_single_slide_copilot(request: SingleSlideCopilotRequest) -> SingleSlideCopilotResult:
    """Apply bounded deterministic edits to one slide and preserve the original."""

    original = request.original
    updated = original.model_copy(deep=True)
    instruction = request.instruction.strip().lower()
    action_items: List[str] = []
    if any(token in instruction for token in ("减少文字", "精简", "shorten", "shorter")):
        updated.paragraph = _shorten(updated.paragraph, max(120, updated.character_budget // 2))
        updated.bullets = updated.bullets[:3]
        action_items.append("已限制段落长度并保留前三条要点。")
    if any(token in instruction for token in ("科研贡献", "贡献", "contribution")):
        updated.bullets = [
            bullet.replace("项目：", "科研贡献：", 1) if bullet.startswith("项目：") else bullet
            for bullet in updated.bullets
        ]
        action_items.append("已将项目要点改写为贡献导向标签。")
    if any(token in instruction for token in ("流程图", "流程", "process")):
        updated.bullets = [
            f"{index}. {bullet}" for index, bullet in enumerate(updated.bullets, start=1)
        ]
        action_items.append("已将要点改为顺序结构，供后续流程图适配器使用。")
    if any(token in instruction for token in ("讲述", "speaker", "口述")):
        updated.speaker_notes.append("讲述时先说明问题，再说明个人贡献，最后连接导师方向。")
        action_items.append("已补充讲述重点。")
    if updated == original:
        action_items.append("当前指令没有匹配到受支持的局部修改动作。")
    issues = _validate_slide(updated)
    diff_text = "\n".join(
        difflib.unified_diff(
            _slide_lines(original),
            _slide_lines(updated),
            fromfile="original-slide",
            tofile="updated-slide",
            lineterm="",
        )
    )
    return SingleSlideCopilotResult(
        slide_index=request.slide_index,
        original=original,
        updated=updated,
        diff_text=diff_text,
        passed=not issues,
        issues=issues,
        action_items=action_items,
    )


def save_single_slide_version(
    workspace,
    request: SingleSlideCopilotRequest,
    result: SingleSlideCopilotResult,
) -> SingleSlideVersionRecord:
    """Persist an auditable local version without mutating the source slide."""

    digest = hashlib.sha256(
        f"{request.slide_index}:{request.instruction}:{result.diff_text}".encode("utf-8")
    ).hexdigest()[:12]
    record = SingleSlideVersionRecord(
        version_id=f"slidever_{digest}",
        slide_index=request.slide_index,
        instruction=request.instruction,
        original=request.original,
        updated=result.updated,
        diff_text=result.diff_text,
        passed=result.passed,
        issues=result.issues,
        action_items=result.action_items,
    )
    dump = record.model_dump() if hasattr(record, "model_dump") else record.dict()
    workspace.write("presentation_slide_versions", dump, "version_id")
    return record


def _element_kind(shape, text_index: int) -> SlideElementKind:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    text = getattr(shape, "text", "").strip()
    if not text:
        return "shape"
    lowered = text.lower()
    if text_index == 0:
        return "title"
    if any(token in text for token in ("日期", "2026", "2025")):
        return "date"
    if any(token in lowered for token in ("logo", "brand")):
        return "logo"
    if any(token in text for token in ("谢谢", "thank", "Q&A")):
        return "footer"
    return "body"


def _shape_description(shape) -> str:
    text = getattr(shape, "text", "").strip().replace("\n", " ")
    if text:
        return text[:120]
    return str(getattr(shape, "shape_type", "unknown"))


def _char_capacity(shape) -> int:
    width = float(getattr(shape, "width", 0)) / 914400
    height = float(getattr(shape, "height", 0)) / 914400
    return max(0, round(width * height * 18))


def _constraints(kind: SlideElementKind) -> List[str]:
    if kind == "title":
        return ["single_line_preferred", "max_120_chars"]
    if kind == "body":
        return ["preserve_evidence_refs", "avoid_unconfirmed_claims"]
    if kind in {"image", "logo"}:
        return ["do_not_fetch_external_asset_implicitly"]
    if kind == "unsupported":
        return ["manual_review_required"]
    return []


def _page_type(slide_index: int, texts: List[str]) -> str:
    joined = " ".join(texts)
    if slide_index == 1:
        return "opening"
    if any(token in joined for token in ("谢谢", "Thank", "Q&A")):
        return "ending"
    if any(token in joined for token in ("目录", "Agenda", "Contents")):
        return "toc"
    return "content"


def _shorten(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    shortened = re.split(r"(?<=[。！？!?；;])", text[:limit])
    return "".join(shortened[:-1] or shortened).strip()


def _slide_lines(slide: OrganizedSlideContent) -> List[str]:
    return [
        f"paragraph: {slide.paragraph}",
        *[f"bullet: {bullet}" for bullet in slide.bullets],
        *[f"note: {note}" for note in slide.speaker_notes],
        f"evidence: {', '.join(slide.evidence_refs)}",
    ]


def _validate_slide(slide: OrganizedSlideContent) -> List[str]:
    issues = []
    if len(slide.paragraph) > slide.character_budget:
        issues.append("paragraph_exceeds_character_budget")
    if any(not item.strip() for item in slide.bullets):
        issues.append("empty_bullet")
    if set(slide.rejected_fields) & set(slide.unconfirmed_fields):
        issues.append("field_state_conflict")
    return issues
