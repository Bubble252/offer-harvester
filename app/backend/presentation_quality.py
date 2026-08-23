from __future__ import annotations

import hashlib
import re
from pathlib import Path

from models import (
    GeneratedMaterial,
    PresentationGenerationRequest,
    PresentationPrecheckIssue,
    PresentationPrecheckReport,
    PresentationQualityReport,
    PresentationTaskRecord,
    ReferencePresentationRecord,
)
from pptx import Presentation
from storage import Workspace, safe_filename

from integrations.presentation_engine import PresentationResult

ALLOWED_REFERENCE_SUFFIXES = {".pptx"}


def save_reference_presentation(
    workspace: Workspace,
    content: bytes,
    original_filename: str,
) -> tuple[ReferencePresentationRecord, PresentationPrecheckReport]:
    if not content:
        raise ValueError("Reference PPT content is required")
    filename = safe_filename(original_filename or "reference.pptx")
    suffix = Path(filename).suffix.lower()
    if suffix not in ALLOWED_REFERENCE_SUFFIXES:
        raise ValueError("Only .pptx reference files are supported in the first version")

    digest = hashlib.sha256(content).hexdigest()
    reference_id = f"ppt_ref_{digest[:12]}"
    relative_path = Path("reference_presentations") / f"{reference_id}_{filename}"
    path = workspace.root / relative_path
    path.write_bytes(content)
    record = ReferencePresentationRecord(
        reference_id=reference_id,
        original_filename=original_filename or filename,
        path=relative_path.as_posix(),
        content_hash=f"sha256:{digest}",
        notes="参考 PPT 只用于预检和后续可选 adapter，不直接进入事实来源。",
    )
    workspace.write("reference_presentations", dump(record), "reference_id")
    precheck = precheck_reference_presentation(path, record.reference_id, record.original_filename)
    workspace.write("presentation_prechecks", dump(precheck), "precheck_id")
    return record, precheck


def precheck_reference_presentation(
    path: Path,
    reference_id: str = "",
    filename: str = "",
) -> PresentationPrecheckReport:
    issues: list[PresentationPrecheckIssue] = []
    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        return PresentationPrecheckReport(
            reference_id=reference_id,
            filename=filename or path.name,
            passed=False,
            issues=[
                PresentationPrecheckIssue(
                    code="pptx_parse_failed",
                    message=f"参考 PPT 无法解析：{exc}",
                    severity="error",
                )
            ],
        )

    shape_counts = [len(slide.shapes) for slide in presentation.slides]
    functional_pages = detect_functional_pages(presentation)
    for index, count in enumerate(shape_counts, start=1):
        if count > 45:
            issues.append(
                PresentationPrecheckIssue(
                    code="shape_count_high",
                    message="单页元素数量较多，外部模板适配时可能失败。",
                    severity="warning",
                    slide_index=index,
                )
            )
    if len(presentation.slides) == 0:
        issues.append(
            PresentationPrecheckIssue(
                code="empty_deck",
                message="参考 PPT 没有页面。",
                severity="error",
            )
        )
    if not functional_pages.get("opening"):
        issues.append(
            PresentationPrecheckIssue(
                code="opening_missing",
                message="未检测到明显封面页，后续可降级使用本地模板。",
                severity="info",
            )
        )
    return PresentationPrecheckReport(
        reference_id=reference_id,
        filename=filename or path.name,
        passed=not any(issue.severity == "error" for issue in issues),
        slide_count=len(presentation.slides),
        total_shape_count=sum(shape_counts),
        max_shapes_per_slide=max(shape_counts or [0]),
        functional_pages=functional_pages,
        issues=issues,
        fallback_allowed=True,
    )


def build_presentation_quality_report(
    task: PresentationTaskRecord,
    outline: GeneratedMaterial,
    result: PresentationResult,
    request: PresentationGenerationRequest,
) -> PresentationQualityReport:
    content_score = score_content(outline)
    design_score = 75
    coherence_score = score_coherence(outline, request)
    issues: list[PresentationPrecheckIssue] = []
    action_items: list[str] = []

    if result.output_path and result.output_path.exists():
        try:
            presentation = Presentation(str(result.output_path))
            shape_counts = [len(slide.shapes) for slide in presentation.slides]
            if len(presentation.slides) != max(1, min(request.num_slides, 12)):
                issues.append(
                    PresentationPrecheckIssue(
                        code="slide_count_mismatch",
                        message="生成页数与请求页数不一致。",
                        severity="warning",
                    )
                )
                action_items.append("复核 PPT 大纲章节数和目标页数。")
            if max(shape_counts or [0]) > 35:
                design_score -= 10
                issues.append(
                    PresentationPrecheckIssue(
                        code="generated_shape_count_high",
                        message="生成 PPT 单页元素较多，可能影响可读性。",
                        severity="warning",
                    )
                )
        except Exception as exc:
            design_score -= 20
            issues.append(
                PresentationPrecheckIssue(
                    code="generated_parse_failed",
                    message=f"生成 PPT 无法复检：{exc}",
                    severity="warning",
                )
            )
    else:
        design_score -= 30
        action_items.append("重新生成 PPTX 或下载 Markdown 大纲兜底。")

    if content_score < 70:
        action_items.append("补充学生项目、导师方向和未来计划的具体证据。")
    if coherence_score < 70:
        action_items.append("按封面、背景、项目、匹配、计划顺序重排大纲。")

    total = round(content_score * 0.45 + design_score * 0.3 + coherence_score * 0.25)
    return PresentationQualityReport(
        task_id=task.task_id,
        target_id=task.target_id,
        outline_material_id=outline.material_id,
        engine_name=result.engine_name,
        content_score=max(0, min(100, content_score)),
        design_score=max(0, min(100, design_score)),
        coherence_score=max(0, min(100, coherence_score)),
        total_score=max(0, min(100, total)),
        issues=issues,
        action_items=action_items,
    )


def detect_functional_pages(presentation: Presentation) -> dict[str, bool]:
    texts = [slide_text(slide).lower() for slide in presentation.slides]
    return {
        "opening": bool(texts and len(texts[0]) > 0),
        "toc": any("目录" in text or "agenda" in text or "contents" in text for text in texts),
        "section_header": any("section" in text or "章节" in text for text in texts),
        "ending": any("谢谢" in text or "thank" in text or "q&a" in text for text in texts),
    }


def slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            parts.append(shape.text)
    return "\n".join(parts)


def score_content(outline: GeneratedMaterial) -> int:
    text = outline.content
    score = 55
    for token in ["项目", "研究", "导师", "方向", "计划"]:
        if token in text:
            score += 7
    if len(re.findall(r"^- ", text, flags=re.M)) >= 8:
        score += 10
    return min(score, 100)


def score_coherence(outline: GeneratedMaterial, request: PresentationGenerationRequest) -> int:
    headings = re.findall(r"^##\s+(.+)$", outline.content, flags=re.M)
    score = 55
    if len(headings) >= min(request.num_slides, 5):
        score += 20
    joined = " ".join(headings)
    for token in ["封面", "教育", "项目", "匹配", "计划"]:
        if token in joined:
            score += 5
    return min(score, 100)


def dump(model):
    return model.model_dump() if hasattr(model, "model_dump") else model.dict()
