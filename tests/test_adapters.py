import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "app" / "backend"))

from models import AdvisorProfile, RAGSearchHit, StudentProfile  # noqa: E402
from presentation_quality import precheck_reference_presentation  # noqa: E402
from presentation_workflow import (  # noqa: E402
    SingleSlideCopilotRequest,
    build_presentation_schema,
    organize_profile_content,
    run_single_slide_copilot,
    save_single_slide_version,
)
from storage import Workspace  # noqa: E402

from integrations.presentation_engine.adapter import (  # noqa: E402
    LocalOutlineAdapter,
    LocalPptxAdapter,
    PresentationRequest,
)
from integrations.workflow_engine.adapter import WorkflowLogger  # noqa: E402


def test_local_presentation_adapter_has_a_safe_fallback(tmp_path):
    result = LocalOutlineAdapter().generate(
        PresentationRequest(
            title="张三/保研面试",
            outline="# 面试展示\n",
            output_dir=tmp_path,
        )
    )

    assert result.status == "fallback"
    assert result.output_path is not None
    assert result.output_path.exists()
    assert result.output_path.read_text(encoding="utf-8") == "# 面试展示\n"
    assert "/" not in result.output_path.name


def test_workflow_logger_keeps_structured_events():
    logger = WorkflowLogger()
    event = logger.record("profile_created", {"profile_id": "profile_demo"})

    assert event.event_type == "profile_created"
    assert event.payload["profile_id"] == "profile_demo"
    assert logger.events == [event]


def test_local_pptx_adapter_generates_five_editable_slides(tmp_path):
    outline = """# 5 分钟保研面试展示 PPT 大纲

## 1. 封面
- 标题：匿名学生

## 2. 教育背景
- 学校：某大学

## 3. 项目经历
- 项目：多模态问答

## 4. 方向匹配
- 导师方向：多模态学习

## 5. 未来计划
- 阅读课题组论文
"""
    result = LocalPptxAdapter().generate(
        PresentationRequest(title="匿名学生面试展示", outline=outline, output_dir=tmp_path)
    )

    assert result.status == "success"
    assert result.output_path is not None
    assert result.output_path.suffix == ".pptx"
    assert result.output_path.exists()

    from pptx import Presentation

    presentation = Presentation(str(result.output_path))
    assert len(presentation.slides) == 5
    assert presentation.slide_width > presentation.slide_height


def test_local_pptx_adapter_respects_requested_slide_count(tmp_path):
    outline = "\n".join(
        [
            "# 面试展示",
            *[f"\n## {index}. 页面 {index}\n- 内容 {index}\n" for index in range(1, 8)],
        ]
    )
    result = LocalPptxAdapter().generate(
        PresentationRequest(
            title="三页面试展示",
            outline=outline,
            output_dir=tmp_path,
            num_slides=3,
        )
    )

    from pptx import Presentation

    presentation = Presentation(str(result.output_path))
    assert len(presentation.slides) == 3


def test_reference_ppt_precheck_reports_functional_pages(tmp_path):
    result = LocalPptxAdapter().generate(
        PresentationRequest(
            title="参考模板",
            outline="# 参考模板\n\n## 1. 封面\n- 标题\n\n## 2. Agenda\n- 目录\n",
            output_dir=tmp_path,
            num_slides=2,
        )
    )

    report = precheck_reference_presentation(
        result.output_path,
        reference_id="ppt_ref_demo",
        filename="reference.pptx",
    )

    assert report.passed is True
    assert report.slide_count == 2
    assert report.functional_pages["opening"] is True
    assert report.fallback_allowed is True


def test_presentation_schema_extracts_slots_and_unsupported_elements(tmp_path):
    result = LocalPptxAdapter().generate(
        PresentationRequest(
            title="参考模板",
            outline=(
                "# 参考模板\n\n## 1. 封面\n- 标题：匿名学生\n\n## 2. 项目\n- 项目：证据审计\n"
            ),
            output_dir=tmp_path,
            num_slides=2,
        )
    )

    schema = build_presentation_schema(result.output_path, source_reference_id="ppt_ref_demo")
    assert schema.schema_version == "17c-v1"
    assert schema.source_path == "ppt_ref_demo"
    assert len(schema.slides) == 2
    assert schema.slides[0].page_type == "opening"
    assert any(element.kind == "title" for element in schema.slides[0].elements)
    assert all(element.char_capacity >= 0 for slide in schema.slides for element in slide.elements)


def test_content_organizer_and_single_slide_copilot_preserve_evidence_and_scope(tmp_path):
    profile = StudentProfile(
        education="某大学计算机学院",
        projects=["证据审计工作流"],
        research_interests=["多模态学习"],
        confirmation_map={"projects": "unconfirmed", "research_interests": "confirmed"},
    )
    advisor = AdvisorProfile(research_directions=["多模态学习"])
    evidence = [
        RAGSearchHit(
            source_id="kb_demo",
            chunk_id="chunk_demo",
            title="导师主页",
            evidence_ref="kb_demo#chunk_demo",
        )
    ]
    organized = organize_profile_content(profile, advisor, evidence)
    slide = organized.slides[0]
    assert slide.bullets
    assert slide.evidence_refs == ["kb_demo#chunk_demo"]
    assert "projects" in slide.unconfirmed_fields

    result = run_single_slide_copilot(
        SingleSlideCopilotRequest(
            slide_index=1,
            instruction="减少文字并突出科研贡献",
            original=slide,
        )
    )
    assert result.passed
    assert result.updated.slide_index == slide.slide_index
    assert result.updated.evidence_refs == slide.evidence_refs
    assert result.diff_text
    assert result.original.bullets != result.updated.bullets

    workspace = Workspace(str(tmp_path))
    version = save_single_slide_version(
        workspace,
        SingleSlideCopilotRequest(slide_index=1, instruction="减少文字", original=slide),
        run_single_slide_copilot(
            SingleSlideCopilotRequest(slide_index=1, instruction="减少文字", original=slide)
        ),
    )
    assert workspace.read("presentation_slide_versions", version.version_id)
