from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional, Protocol

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


@dataclass
class PresentationRequest:
    """演示文稿生成请求，业务层只依赖这一组稳定字段。"""

    title: str
    outline: str
    output_dir: Path
    reference_file: Optional[Path] = None
    metadata: Dict[str, str] = field(default_factory=dict)


@dataclass
class PresentationResult:
    """统一描述生成结果和降级原因。"""

    status: str
    output_path: Optional[Path] = None
    message: str = ""


@dataclass
class PresentationTask:
    """描述一次本地 PPTX 生成任务的持久化状态。"""

    task_id: str
    status: str
    output_path: Optional[Path] = None
    error: str = ""


class PresentationAdapter(Protocol):
    """约束外部演示文稿引擎必须提供的最小能力。"""

    def generate(self, request: PresentationRequest) -> PresentationResult:
        ...


class LocalPptxAdapter:
    """生成紫罗兰主题的原生文本框 PPTX，不依赖外部生成引擎。"""

    PURPLE = RGBColor(91, 64, 145)
    LIGHT_PURPLE = RGBColor(242, 238, 250)
    INK = RGBColor(39, 35, 52)
    MUTED = RGBColor(100, 94, 114)
    WHITE = RGBColor(255, 255, 255)

    def generate(self, request: PresentationRequest) -> PresentationResult:
        request.output_dir.mkdir(parents=True, exist_ok=True)
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        self._remove_default_slides(presentation)
        sections = parse_outline(request.outline)
        for index, section in enumerate(sections[:5], start=1):
            self._add_slide(presentation, request.title, section, index, len(sections[:5]))
        output_path = request.output_dir / f"{safe_name(request.title)}.pptx"
        presentation.save(str(output_path))
        return PresentationResult(
            status="success",
            output_path=output_path,
            message="已生成可编辑 PPTX。",
        )

    @staticmethod
    def _remove_default_slides(presentation: Presentation) -> None:
        while presentation.slides:
            slide_id = presentation.slides._sldIdLst[0]
            presentation.part.drop_rel(slide_id.rId)
            presentation.slides._sldIdLst.remove(slide_id)

    def _add_slide(
        self,
        presentation: Presentation,
        deck_title: str,
        section: Dict[str, object],
        index: int,
        total: int,
    ) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = self.WHITE

        accent = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(0.22), Inches(7.5)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.PURPLE
        accent.line.fill.background()

        top_rule = slide.shapes.add_shape(
            1, Inches(0.75), Inches(0.62), Inches(11.85), Inches(0.04)
        )
        top_rule.fill.solid()
        top_rule.fill.fore_color.rgb = self.LIGHT_PURPLE
        top_rule.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.78), Inches(0.82), Inches(10.9), Inches(0.65)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        paragraph = title_frame.paragraphs[0]
        paragraph.text = str(section["title"])
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(25 if index == 1 else 22)
        paragraph.font.bold = True
        paragraph.font.color.rgb = self.INK

        if index == 1:
            subtitle = slide.shapes.add_textbox(
                Inches(0.82), Inches(1.65), Inches(10.4), Inches(0.6)
            )
            subtitle.text_frame.text = deck_title
            subtitle.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
            subtitle.text_frame.paragraphs[0].font.size = Pt(17)
            subtitle.text_frame.paragraphs[0].font.color.rgb = self.PURPLE

        bullets: List[str] = section["bullets"]  # type: ignore[assignment]
        body_box = slide.shapes.add_textbox(
            Inches(0.95), Inches(2.0 if index == 1 else 1.72), Inches(10.9), Inches(4.7)
        )
        body = body_box.text_frame
        body.clear()
        body.word_wrap = True
        for bullet_index, bullet in enumerate(bullets[:7]):
            paragraph = body.paragraphs[0] if bullet_index == 0 else body.add_paragraph()
            paragraph.text = f"• {bullet}"
            paragraph.level = 0
            paragraph.font.name = "Microsoft YaHei"
            paragraph.font.size = Pt(18 if index == 1 else 16)
            paragraph.font.color.rgb = self.INK
            paragraph.space_after = Pt(13)

        footer = slide.shapes.add_textbox(
            Inches(0.82), Inches(6.92), Inches(11.3), Inches(0.28)
        )
        footer_frame = footer.text_frame
        footer_frame.clear()
        footer_paragraph = footer_frame.paragraphs[0]
        footer_paragraph.text = f"保研硕博申请展示  ·  {index:02d} / {total:02d}"
        footer_paragraph.font.name = "Microsoft YaHei"
        footer_paragraph.font.size = Pt(9)
        footer_paragraph.font.color.rgb = self.MUTED
        footer_paragraph.alignment = PP_ALIGN.RIGHT


class LocalOutlineAdapter:
    """在未配置外部引擎时保留可下载的 Markdown 大纲。"""

    def generate(self, request: PresentationRequest) -> PresentationResult:
        request.output_dir.mkdir(parents=True, exist_ok=True)
        output_path = request.output_dir / f"{safe_name(request.title)}.md"
        output_path.write_text(request.outline, encoding="utf-8")
        return PresentationResult(
            status="fallback",
            output_path=output_path,
            message="未配置演示文稿引擎，已生成可审阅的 Markdown 大纲。",
        )


def parse_outline(outline: str) -> List[Dict[str, object]]:
    """Convert the stable Markdown outline into slide-sized structured content."""

    sections: List[Dict[str, object]] = []
    current: Optional[Dict[str, object]] = None
    for raw_line in outline.splitlines():
        line = raw_line.strip()
        if line.startswith("## "):
            if current:
                sections.append(current)
            current = {"title": line[3:].strip(), "bullets": []}
        elif line.startswith("- ") and current:
            current["bullets"].append(line[2:].strip())  # type: ignore[union-attr]
    if current:
        sections.append(current)
    return sections


def safe_name(value: str) -> str:
    """限制文件名字符，避免用户输入影响工作区路径。"""

    cleaned = "".join(char if char.isalnum() or char in "-_" else "_" for char in value)
    return cleaned.strip("_")[:80] or "presentation"
